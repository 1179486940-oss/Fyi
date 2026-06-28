"""
MCP: PPT 生成接口
封装 BA 提供的接口2 —— PPT 生成
"""

from __future__ import annotations

from typing import Any, Optional

from config import get_settings
from utils.logger import get_logger
from utils.helpers import generate_uuid

logger = get_logger(__name__)


class PPTGenerateMCP:
    """PPT 生成 MCP Server
    封装 BA 提供的接口2，负责根据数据和模板生成 PPT
    """

    def __init__(self):
        settings = get_settings()
        self._storage_path = settings.storage.storage_path / "ppts"
        self._storage_path.mkdir(parents=True, exist_ok=True)
        self._base_url = f"http://{settings.server.host}:{settings.server.port}/files"

        # TODO: 配置 BA 接口2 的实际地址
        self._api_base_url = ""
        self._api_key = ""

    async def generate_ppt(
        self,
        template_id: str,
        data: dict[str, Any],
        ppt_config: Optional[dict[str, Any]] = None,
    ) -> dict[str, Any]:
        """
        调用 BA 接口2 生成 PPT
        Args:
            template_id: 模板ID
            data: PPT 数据内容
            ppt_config: PPT 配置（标题、作者、页数等）
        Returns:
            生成结果，含下载链接
        """
        logger.info("PPT generate MCP: template=%s", template_id)

        # TODO: 对接 BA 真实接口
        # 当前使用 python-pptx 生成简单的 Mock PPT
        try:
            file_id = generate_uuid()
            file_path = self._storage_path / f"{file_id}.pptx"

            await self._generate_mock_ppt(
                file_path=str(file_path),
                template_id=template_id,
                data=data,
                config=ppt_config or {},
            )

            return {
                "status": "success",
                "file_id": file_id,
                "download_url": f"{self._base_url}/{file_id}/download",
                "format": "pptx",
                "template_used": template_id,
            }

        except Exception as e:
            logger.error("PPT generation failed: %s", e)
            return {"status": "error", "message": str(e)}

    async def _generate_mock_ppt(
        self,
        file_path: str,
        template_id: str,
        data: dict[str, Any],
        config: dict[str, Any],
    ) -> None:
        """使用 python-pptx 生成本地 Mock PPT"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ── 封面 ─────────────────────────────────────
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # 标题
        left, top, width, height = Inches(1), Inches(2.5), Inches(11), Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = config.get("title", "PPT 演示文稿")
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        left, top, width, height = Inches(1), Inches(4.2), Inches(11), Inches(1)
        txBox2 = slide.shapes.add_textbox(left, top, width, height)
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = config.get("subtitle", f"模板: {template_id}")
        p2.font.size = Pt(20)
        p2.alignment = PP_ALIGN.CENTER

        # ── 内容页（从 data 构造）─────────────────────
        slides_data = data.get("slides", data.get("data", []))
        if isinstance(slides_data, list):
            for item in slides_data[:10]:  # 最多10页
                slide = prs.slides.add_slide(slide_layout)
                left, top, width, height = Inches(1), Inches(1.5), Inches(11), Inches(4)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = str(item.get("title", item.get("content", "")))
                p.font.size = Pt(28)
        else:
            # 单页内容
            slide = prs.slides.add_slide(slide_layout)
            left, top, width, height = Inches(1), Inches(1.5), Inches(11), Inches(4)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = str(slides_data)
            p.font.size = Pt(24)

        # ── 结尾页 ─────────────────────────────────────
        slide = prs.slides.add_slide(slide_layout)
        left, top, width, height = Inches(1), Inches(3), Inches(11), Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "感谢观看"
        p.font.size = Pt(40)
        p.alignment = PP_ALIGN.CENTER

        prs.save(file_path)
        logger.info("Mock PPT saved: %s (%d slides)", file_path, len(prs.slides))
